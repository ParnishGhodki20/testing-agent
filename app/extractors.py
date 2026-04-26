"""
File-to-text extractors.

Each extractor receives an absolute file path and returns a plain-text string.
Unsupported extensions return "" and log a warning so the UI can inform the user.
"""
import asyncio
import logging
from pathlib import Path

import pandas as pd
import pypdf
from docx import Document as DocxDocument

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Async extractors
# ---------------------------------------------------------------------------

async def _process_pdf(path: str) -> str:
    parts = []
    reader = pypdf.PdfReader(path)
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


async def _process_docx(path: str) -> str:
    doc = DocxDocument(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


async def _process_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        logger.warning("python-pptx not installed — cannot extract PPTX. Run: pip install python-pptx")
        return ""
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Sync extractors (run via asyncio.to_thread in extract_text)
# ---------------------------------------------------------------------------

def _process_xlsx(path: str) -> str:
    try:
        df = pd.read_excel(path, sheet_name=None)  # all sheets
        parts = []
        for sheet_name, sheet_df in df.items():
            parts.append(f"[Sheet: {sheet_name}]")
            parts.append(sheet_df.to_string(index=False))
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("XLSX extraction failed: %s", exc)
        return ""


def _process_csv(path: str) -> str:
    try:
        return pd.read_csv(path).to_string(index=False)
    except Exception as exc:
        logger.warning("CSV extraction failed: %s", exc)
        return ""


def _process_text(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


# ---------------------------------------------------------------------------
# Dispatch tables
# ---------------------------------------------------------------------------

_ASYNC_PROCESSORS: dict = {
    "pdf":  _process_pdf,
    "docx": _process_docx,
    "pptx": _process_pptx,
}

_SYNC_PROCESSORS: dict = {
    "xlsx":    _process_xlsx,
    "csv":     _process_csv,
    "txt":     _process_text,
    "feature": _process_text,
    "json":    _process_text,
    "xml":     _process_text,
    "html":    _process_text,
    "sql":     _process_text,
    "ts":      _process_text,
    "js":      _process_text,
    "java":    _process_text,
    "py":      _process_text,
    "cs":      _process_text,
    "ps1":     _process_text,
    "md":      _process_text,
    "eml":     _process_text,
}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

async def extract_text(file_path: str, extension: str) -> str:
    """
    Extract plain text from a file.

    Args:
        file_path: Absolute path to the file.
        extension: File extension without leading dot (e.g. "pdf", "docx").

    Returns:
        Extracted text string, or "" if unsupported / error.
    """
    ext = extension.lstrip(".").lower()

    if ext in _ASYNC_PROCESSORS:
        try:
            return await _ASYNC_PROCESSORS[ext](file_path)
        except Exception as exc:
            logger.error("Async extractor failed for .%s: %s", ext, exc)
            return ""

    if ext in _SYNC_PROCESSORS:
        try:
            return await asyncio.to_thread(_SYNC_PROCESSORS[ext], file_path)
        except Exception as exc:
            logger.error("Sync extractor failed for .%s: %s", ext, exc)
            return ""

    # Unsupported — warn explicitly so caller can surface message to user
    logger.warning("No extractor for extension '.%s' — file will be skipped.", ext)
    return ""